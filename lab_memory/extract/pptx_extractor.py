"""PPTX text extractor for weekly lab reports."""
from __future__ import annotations

import json
import os
import re
from datetime import datetime
from multiprocessing import Pool, cpu_count
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.table import Table


def _table_to_markdown(table: Table) -> str:
    """Convert a PPTX table to markdown format."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")
    if len(rows) >= 1:
        # Add header separator after first row
        header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
        rows.insert(1, header_sep)
    return "\n".join(rows)


def _parse_date_from_filename(filename: str) -> str | None:
    """Try to extract a date from the filename."""
    patterns = [
        (r"(\d{4})-(\d{2})-(\d{2})", "%Y-%m-%d"),
        (r"(\d{4})(\d{2})(\d{2})", "%Y%m%d"),
        (r"(\d{2})\.(\d{2})\.(\d{2})", "%y.%m.%d"),
        (r"(\d{4})\.(\d{2})\.(\d{2})", "%Y.%m.%d"),
    ]
    for pattern, fmt in patterns:
        m = re.search(pattern, filename)
        if m:
            try:
                date_str = m.group(0)
                dt = datetime.strptime(date_str, fmt)
                return dt.strftime("%Y-%m-%d")
            except ValueError:
                continue
    return None


def extract_slide(slide, slide_number: int) -> dict[str, Any]:
    """Extract content from a single slide."""
    title = ""
    body_parts: list[str] = []
    tables: list[str] = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            if slide.shapes.title is not None and shape == slide.shapes.title:
                title = shape.text_frame.text.strip()
            else:
                text = shape.text_frame.text.strip()
                if text:
                    body_parts.append(text)

        if shape.has_table:
            tables.append(_table_to_markdown(shape.table))

    notes = ""
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()

    return {
        "slide_number": slide_number,
        "title": title,
        "body": "\n\n".join(body_parts),
        "tables": tables,
        "notes": notes,
    }


def extract_pptx(file_path: Path) -> dict[str, Any]:
    """Extract all content from a PPTX file."""
    prs = Presentation(str(file_path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        slides.append(extract_slide(slide, i))

    date = _parse_date_from_filename(file_path.stem)

    return {
        "source_file": file_path.name,
        "source_path": str(file_path),
        "date": date,
        "type": "pptx",
        "total_slides": len(slides),
        "slides": slides,
    }


def _extract_single(args: tuple[Path, Path]) -> str:
    """Worker function for multiprocessing."""
    file_path, output_dir = args
    try:
        result = extract_pptx(file_path)
        # Use date or filename as output name
        out_name = result["date"] or file_path.stem
        out_path = output_dir / f"{out_name}.json"
        # Handle duplicates
        counter = 1
        while out_path.exists():
            out_path = output_dir / f"{out_name}_{counter}.json"
            counter += 1
        out_path.write_text(
            json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        return str(out_path)
    except Exception as e:
        return f"ERROR:{file_path}:{e}"


def extract_all(
    input_dir: Path,
    output_dir: Path,
    workers: int | None = None,
    exclude_patterns: list[str] | None = None,
) -> list[str]:
    """Extract all PPTX files from input_dir using multiprocessing."""
    output_dir.mkdir(parents=True, exist_ok=True)
    exclude = [p.lower() for p in (exclude_patterns or [])]
    pptx_files = sorted(
        Path(os.path.join(root, f))
        for root, _, files in os.walk(input_dir, followlinks=True)
        if not any(ex in root.lower() for ex in exclude)
        for f in files
        if f.endswith(".pptx")
    )

    if not pptx_files:
        return []

    args = [(f, output_dir) for f in pptx_files]
    n_workers = workers or min(cpu_count(), len(pptx_files))

    with Pool(n_workers) as pool:
        results = pool.map(_extract_single, args)

    return results
